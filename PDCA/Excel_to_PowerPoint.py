import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

import win32com.client as win32
from win32com.client import constants as pjconstants
import pandas as pd
import pytz
from datetime import datetime
import tkinter as tk
from tkinter import filedialog, simpledialog, messagebox
import math
from openpyxl import load_workbook
import re
import sys
import os
from pyxll import xl_func, xl_macro
############################
POWER_POINT_PATH = None
DATA_FRAME = None
DEMO = True
START_DATE = None
END_DATE = None
PROJEKT_NAME = None
WORKDAYS = None
TODAY = None
############################
def choose_power_point_file():
    global POWER_POINT_PATH
    root = tk.Tk()
    root.withdraw()
    
    POWER_POINT_PATH = filedialog.askopenfilename(filetypes=[("PowerPoint Files", "*.pptx")])

def calculate_workdays(start_date, end_date, holidays=[]):
    global WORKDAYS
    date_range = pd.date_range(start=start_date, end=end_date,freq='B')
    workdays = date_range[date_range.weekday < 5]
    
    for holiday in holidays:
        if holiday in workdays:
            workdays = workdays.drop(holiday)
    
    WORKDAYS =  len(workdays)
    return len(workdays)
    
def main(excel_file_path):
    global POWER_POINT_PATH
    global DATA_FRAME
    global START_DATE
    global END_DATE
    global WORKDAYS
    global TODAY
    global PROJEKT_NAME
    
    TODAY = datetime.now().strftime("%Y-%m-%d")

    
    print(TODAY)
    
    workbook = load_workbook(excel_file_path)
    worksheet = workbook['Cockpit']
    
    DATA_FRAME = pd.DataFrame(worksheet.values)
    DATA_FRAME = DATA_FRAME.reset_index()
    PROJEKT_NAME = DATA_FRAME.at[2,3]
    
    for _,row in DATA_FRAME.iterrows():
        if row[3] == "Start":
            START_DATE = row[5]
        elif row[3] == "Ende":
            END_DATE = row[5]
            break
        
        
    planned = calculate_workdays(START_DATE, END_DATE, [])
    
    status = calculate_workdays(TODAY,END_DATE,[])
    
    rest = planned - status
       
    choose_power_point_file()
    presentation = Presentation(POWER_POINT_PATH)

    labels = ['Ist', 'Rest']
    
    sizes = [status/planned, rest/planned]

    colors = ['green', 'red']
    
    explode = (0,0.1)

    plt.pie(sizes,labels=labels,colors=colors,autopct='%1.1f%%',startangle=140,explode=explode)
    
    plt.title(PROJEKT_NAME)
    plt.savefig(f'{PROJEKT_NAME}.png', format= 'png')
    plt.show()


    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    left = Inches(1)
    top = Inches(1)
    width = Inches(9)
    height = Inches(6)


    slide.shapes.add_picture(f'{PROJEKT_NAME}.png',left,top,width,height)

    presentation.save(r'C:\Users\npawelka\Desktop\Example.pptx')

if __name__ == "__main__":
    if DEMO is False:
        if len(sys.argv) != 2:
            messagebox.showerror("Error","Fehler bei der Ermittlung des Commands")
            sys.exit()
        else:
            main(sys.argv[1])
    else:
        main(r"C:\Users\npawelka\Desktop\PDCA\ETO.EEVACTUATOR.Entw.016.xlsm")
